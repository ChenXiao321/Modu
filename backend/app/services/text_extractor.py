import os
from pathlib import Path


class TextExtractorError(Exception):
    """Raised when text extraction fails."""

    def __init__(self, message: str) -> None:
        self.message = message
        super().__init__(message)


class TextExtractor:
    """Extract plain text from supported document formats."""

    _MIME_MAP: dict[str, str] = {
        "application/pdf": ".pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        "text/plain": ".txt",
    }

    @staticmethod
    def extract(file_path: str, file_type: str) -> str:
        if not os.path.exists(file_path):
            raise TextExtractorError(f"文件不存在: {file_path}")

        ext = TextExtractor._resolve_ext(file_path, file_type)
        text = ""

        if ext == ".pdf":
            text = TextExtractor._extract_pdf(file_path)
        elif ext in (".docx", ".doc"):
            text = TextExtractor._extract_docx(file_path)
        elif ext == ".xlsx":
            text = TextExtractor._extract_xlsx(file_path)
        elif ext == ".pptx":
            text = TextExtractor._extract_pptx(file_path)
        elif ext == ".txt":
            text = TextExtractor._extract_txt(file_path)
        else:
            raise TextExtractorError(f"不支持的文本提取格式: {ext}")

        if not text.strip():
            raise TextExtractorError("文档未提取到任何文本内容（可能是扫描件、图片或空文件）")

        return text

    @staticmethod
    def _resolve_ext(file_path: str, file_type: str) -> str:
        mapped = TextExtractor._MIME_MAP.get(file_type.lower())
        if mapped:
            return mapped
        ext = Path(file_path).suffix.lower()
        if ext:
            return ext
        raise TextExtractorError(f"无法确定文件类型: file_type={file_type}, path={file_path}")

    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        import pdfplumber

        text_parts: list[str] = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
        except Exception as e:
            raise TextExtractorError(f"PDF 解析失败: {e}") from e
        return "\n\n".join(text_parts)

    @staticmethod
    def _extract_docx(file_path: str) -> str:
        from docx import Document

        try:
            doc = Document(file_path)
        except Exception as e:
            raise TextExtractorError(f"Word 文档解析失败（文件可能损坏或格式不正确）: {e}") from e

        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    @staticmethod
    def _extract_xlsx(file_path: str) -> str:
        from openpyxl import load_workbook

        try:
            wb = load_workbook(file_path, data_only=True)
        except Exception as e:
            raise TextExtractorError(f"Excel 解析失败: {e}") from e

        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    sheet_texts.append(row_text)
            if sheet_texts:
                parts.append(f"[{sheet_name}]\n" + "\n".join(sheet_texts))
        return "\n\n".join(parts)

    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        from pptx import Presentation

        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise TextExtractorError(f"PPT 解析失败: {e}") from e

        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"[Slide {idx}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)

    @staticmethod
    def _extract_txt(file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
